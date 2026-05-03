import os
import uuid
import tempfile
import zipfile
import io
import subprocess
from xml.etree import ElementTree as ET
import json
import fitz
import openpyxl
import xlrd
from fastapi import APIRouter, File, UploadFile, HTTPException, Depends
from pydantic import BaseModel
from pillow_heif import register_heif_opener
from PIL import Image, ImageOps

register_heif_opener()
from typing import List
from docx import Document
from pptx import Presentation
from google.cloud import speech
from .auth import User, get_current_user
from logging_util import logger

router = APIRouter()

IMAGE_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "uploads", "images")
FILES_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "uploads", "files")
FILES_ORIGINAL_DIR = os.path.join(FILES_DIR, "original")
FILES_PROCESSED_DIR = os.path.join(FILES_DIR, "processed")

os.makedirs(IMAGE_DIR, exist_ok=True)
os.makedirs(FILES_DIR, exist_ok=True)
os.makedirs(FILES_ORIGINAL_DIR, exist_ok=True)
os.makedirs(FILES_PROCESSED_DIR, exist_ok=True)

BINARY_SIGNATURES = [
    b'\x89PNG', b'\xff\xd8\xff', b'GIF8', b'RIFF',
    b'MZ', b'\x7fELF', b'\xca\xfe\xba\xbe',
    b'PK\x03\x04', b'\x50\x4b\x03\x04', b'Rar!',
]
TEXT_ENCODINGS = ['utf-8-sig', 'utf-16', 'utf-8', 'cp949', 'euc-kr']

class WebContent(BaseModel):
    unique_id: str
    html: str
    stylesheets: List[str]
    title: str

def extract_text(data: bytes, ext: str, filename: str) -> str:
    try:
        if ext == '.pdf':
            with fitz.open(stream=data, filetype="pdf") as doc:
                return "\n".join(page.get_text() for page in doc).strip()

        if ext == '.docx':
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs).strip()

        if ext == '.xlsx':
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    line = "\t".join(str(c) if c is not None else "" for c in row)
                    if line.strip():
                        parts.append(line)
            return "\n".join(parts).strip()

        if ext == '.xls':
            wb = xlrd.open_workbook(file_contents=data)
            parts = []
            for sheet in wb.sheets():
                for i in range(sheet.nrows):
                    line = "\t".join(str(sheet.cell_value(i, j)) for j in range(sheet.ncols))
                    if line.strip():
                        parts.append(line)
            return "\n".join(parts).strip()

        if ext == '.pptx':
            prs = Presentation(io.BytesIO(data))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return "\n".join(parts).strip()

        if ext == '.hwp':
            with tempfile.NamedTemporaryFile(suffix='.hwp', delete=False) as tmp:
                tmp.write(data)
                tmp_path = tmp.name
            try:
                hwp5txt = os.path.join(os.path.dirname(os.path.dirname(__file__)), ".venv", "bin", "hwp5txt")
                result = subprocess.run([hwp5txt, tmp_path], capture_output=True, text=True, timeout=30)
                if result.returncode != 0:
                    raise Exception(result.stderr)
                return result.stdout.strip()
            finally:
                os.remove(tmp_path)

        if ext == '.hwpx':
            with zipfile.ZipFile(io.BytesIO(data)) as z:
                sections = sorted(n for n in z.namelist() if n.startswith("Contents/section") and n.endswith(".xml"))
                parts = []
                for section in sections:
                    root = ET.fromstring(z.read(section))
                    parts.extend(el.text for el in root.iter() if el.tag.endswith('}t') and 'hancom.co.kr/hwpml' in el.tag and el.text)
            return "\n".join(parts).strip()

        for sig in BINARY_SIGNATURES:
            if data.startswith(sig):
                raise HTTPException(status_code=422, detail=f"Binary file is not supported: {filename}")
        for enc in TEXT_ENCODINGS:
            try:
                return data.decode(enc).strip()
            except (UnicodeDecodeError, Exception):
                continue
        raise HTTPException(status_code=422, detail=f"Binary file is not supported: {filename}")

    except HTTPException:
        raise
    except Exception as ex:
        logger.info(f"EXTRACT_FAILED: {json.dumps({'file': filename, 'error': str(ex)}, ensure_ascii=False)}")
        raise HTTPException(status_code=422, detail="Text extraction failed")

@router.post("/upload/image")
async def upload_image(file: UploadFile = File(...), current_user: User = Depends(get_current_user)):
    file_data = await file.read()
    if not current_user.admin and len(file_data) > 10 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File size exceeds 10MB limit.")

    try:
        image = Image.open(io.BytesIO(file_data))
    except Exception:
        raise HTTPException(status_code=400, detail="Can't read image file.")

    image = ImageOps.exif_transpose(image)

    if image.mode in ("RGBA", "LA") or (image.mode == "P" and "transparency" in image.info):
        if image.mode != "RGBA":
            image = image.convert("RGBA")
        background = Image.new("RGB", image.size, (255, 255, 255))
        background.paste(image, mask=image.split()[3])
        image = background
    else:
        image = image.convert("RGB")

    max_dimension = (1024, 1024)
    image.thumbnail(max_dimension, Image.Resampling.LANCZOS)

    buffer = io.BytesIO()
    image.save(buffer, format="JPEG", quality=70, optimize=True)

    saved_filename = f"{uuid.uuid4().hex}.jpeg"
    file_location = os.path.join(IMAGE_DIR, saved_filename)
    with open(file_location, "wb") as f:
        f.write(buffer.getvalue())

    return {
        "type": "image",
        "name": file.filename,
        "content": f"/uploads/images/{saved_filename}"
    }

@router.post("/upload/file")
async def upload_file(file: UploadFile = File(...), current_user: User = Depends(get_current_user)):
    audio_extensions = ['.wav', '.mp3', '.ogg', '.flac', '.amr', '.amr-wb', '.mulaw', '.alaw', '.webm', '.m4a', '.mp4']

    file_data = await file.read()
    if not current_user.admin and len(file_data) > 10 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File size exceeds 10MB limit.")

    filename = file.filename
    _, ext = os.path.splitext(filename)
    ext = ext.lower()

    extracted_text = ""

    if ext == '.zip':
        extracted_parts = []
        try:
            with zipfile.ZipFile(io.BytesIO(file_data)) as z:
                for inner_filename in z.namelist():
                    if inner_filename.startswith("__MACOSX") or os.path.basename(inner_filename).startswith("._"):
                        continue
                    if inner_filename.endswith('/'):
                        continue

                    inner_ext = os.path.splitext(inner_filename)[1].lower()
                    inner_data = z.read(inner_filename)
                    try:
                        inner_text = extract_text(inner_data, inner_ext, inner_filename)
                    except HTTPException:
                        continue
                    if inner_text.strip():
                        extracted_parts.append(f"[[{inner_filename}]]\n{inner_text}")

            extracted_text = "\n\n".join(extracted_parts)

        except HTTPException:
            raise
        except Exception:
            raise HTTPException(status_code=422, detail="Archive processing failed")

    elif ext in audio_extensions:
        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
            tmp.write(file_data)
            tmp.flush()
            tmp_path = tmp.name

        try:
            client = speech.SpeechClient(client_options={"api_key": os.getenv("GOOGLE_STT_API_KEY")})

            with open(tmp_path, "rb") as audio_file:
                content = audio_file.read()
            audio = speech.RecognitionAudio(content=content)

            config = speech.RecognitionConfig(
                encoding=speech.RecognitionConfig.AudioEncoding.LINEAR16,
                language_code="ko-KR",
                alternative_language_codes=["en-US"],
                enable_automatic_punctuation=True,
            )
            response = client.recognize(config=config, audio=audio)

            for result in response.results:
                extracted_text += result.alternatives[0].transcript + " "
        except Exception as ex:
            logger.info(f"GOOGLE_STT_FAILED: {json.dumps({'file': filename, 'error': str(ex)}, ensure_ascii=False)}")
            extracted_text = extract_text(file_data, ext, filename)
        finally:
            os.remove(tmp_path)

    else:
        extracted_text = extract_text(file_data, ext, filename)

    if not extracted_text.strip():
        raise HTTPException(status_code=422, detail="Text extraction failed")

    if not current_user.admin and len(extracted_text) > 20000:
        raise HTTPException(status_code=413, detail="Extracted text exceeds 20000 character limit.")

    file_uuid = uuid.uuid4().hex

    processed_filename = f"{file_uuid}.txt"
    processed_file_path = os.path.join(FILES_PROCESSED_DIR, processed_filename)
    processed_content = f"[[{filename}]]\n{extracted_text}"
    with open(processed_file_path, "w", encoding="utf-8") as f:
        f.write(processed_content)

    original_filename = f"{file_uuid}{ext}"
    original_file_path = os.path.join(FILES_ORIGINAL_DIR, original_filename)
    with open(original_file_path, "wb") as f:
        f.write(file_data)

    return {
        "type": "file",
        "name": filename,
        "content": f"/uploads/files/processed/{processed_filename}",
        "file_path": f"/uploads/files/original/{original_filename}"
    }

@router.post("/upload_page")
async def upload_page(content: WebContent):
    try:
        stylesheet_content = ""
        for sheet in content.stylesheets:
            if sheet.startswith(('http://', 'https://')):
                stylesheet_content += f'<link rel="stylesheet" href="{sheet}">\n'
            else:
                stylesheet_content += f'{sheet}\n'

        html_content = f"""<!DOCTYPE html>
        <html>
            <head>
                <meta charset="utf-8" />
                <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">
                <meta name="theme-color" content="#000000" />
                <meta name="description" content="DevoChat" />
                <meta property="og:title" content="{content.title}">
                <meta property="og:site_name" content="{content.title}">
                <meta property="og:description" content="DevoChat 공유 링크">
                <meta property="og:image" content="https://devochat.com/full_logo.png">
                <meta property="og:url" content="https://share.devochat.com/og">
                <meta property="og:type" content="website">

                <link rel="icon" type="image/png" href="https://devochat.com/favicon/favicon-96x96.png" sizes="96x96" />
                <link rel="icon" type="image/svg+xml" href="https://devochat.com/favicon/favicon.svg" />
                <link rel="shortcut icon" href="https://devochat.com/favicon/favicon.ico" />
                <link rel="apple-touch-icon" sizes="180x180" href="https://devochat.com/favicon/apple-touch-icon.png" />
                <meta name="apple-mobile-web-app-title" content="DevoChat 공유 링크" />

                <title>{content.title}</title>
                {stylesheet_content}
            </head>
            <body style="display: flex; position: relative; flex-direction: column; margin: 0; overflow: hidden;">
                <div class="header" style="position: sticky; top: 0; background-color: white;">
                    <img src="https://devochat.com/logo.png" alt="DEVOCHAT" width="143.5px" style="cursor: pointer;" onclick="location.href='https://devochat.com'" />
                </div>
                <div class="container">
                    {content.html}
                </div>
            </body>
        </html>"""

        file_path = os.path.join("shared_pages", f"{content.unique_id}.html")
        os.makedirs("shared_pages", exist_ok=True)

        with open(file_path, "w", encoding="utf-8") as f:
            f.write(html_content)

        return {"success": True}

    except Exception as ex:
        raise HTTPException(status_code=500, detail=str(ex))
